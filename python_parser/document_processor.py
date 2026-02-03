"""
Document Processor with Smart Cost Control
Only uses expensive AWS services when really needed
"""

import os
import re
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
import PyPDF2
import pdfplumber
from langchain_text_splitters import RecursiveCharacterTextSplitter
import tiktoken


@dataclass
class ProcessingResult:
    """Result of document processing"""
    success: bool
    text: str
    chunks: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    error: Optional[str] = None
    needs_ocr: bool = False
    needs_vision: bool = False


class DocumentProcessor:
    """
    Smart document processor that:
    1. Tries free text extraction first
    2. Detects if OCR/Vision is needed
    3. Only calls AWS services when necessary
    """
    
    def __init__(self):
        self.encoding = tiktoken.get_encoding("cl100k_base")
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=200,
            length_function=self.count_tokens,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def count_tokens(self, text: str) -> int:
        """Count tokens in text"""
        return len(self.encoding.encode(text))
    
    def process_document(
        self,
        file_path: str,
        file_type: str,
        document_id: str,
        metadata: Dict[str, Any] = None
    ) -> ProcessingResult:
        """
        Process document with smart cost control
        """
        try:
            if file_type == 'pdf':
                return self._process_pdf(file_path, document_id, metadata or {})
            elif file_type == 'docx':
                return self._process_docx(file_path, document_id, metadata or {})
            elif file_type == 'pptx':
                return self._process_pptx(file_path, document_id, metadata or {})
            elif file_type in ['txt', 'html']:
                return self._process_text(file_path, document_id, metadata or {})
            else:
                return ProcessingResult(
                    success=False,
                    text="",
                    chunks=[],
                    metadata={},
                    error=f"Unsupported file type: {file_type}"
                )
        except Exception as e:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error=str(e)
            )
    
    def _process_pdf(
        self,
        file_path: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> ProcessingResult:
        """
        Process PDF with smart extraction
        """
        print(f"📄 Processing PDF: {file_path}")
        
        # Step 1: Try pdfplumber first (best for text PDFs)
        text, quality_score = self._extract_text_pdfplumber(file_path)
        
        # Step 2: Check if text extraction was successful
        needs_ocr = self._needs_ocr(text, quality_score)
        needs_vision = self._needs_vision_analysis(file_path)
        
        if needs_ocr:
            print("⚠️  Low quality text detected - OCR recommended")
            metadata['ocr_recommended'] = True
        
        if needs_vision:
            print("📊 Images/charts detected - Vision analysis recommended")
            metadata['vision_recommended'] = True
        
        # Step 3: Chunk the text
        chunks = self._create_chunks(text, document_id, metadata)
        
        # Step 4: Add processing metadata
        processing_metadata = {
            **metadata,
            'text_length': len(text),
            'token_count': self.count_tokens(text),
            'chunk_count': len(chunks),
            'quality_score': quality_score,
            'extraction_method': 'pdfplumber',
            'needs_ocr': needs_ocr,
            'needs_vision': needs_vision,
        }
        
        return ProcessingResult(
            success=True,
            text=text,
            chunks=chunks,
            metadata=processing_metadata,
            needs_ocr=needs_ocr,
            needs_vision=needs_vision
        )
    
    def _extract_text_pdfplumber(self, file_path: str) -> tuple[str, float]:
        """
        Extract text using pdfplumber
        Returns: (text, quality_score)
        """
        text_parts = []
        total_chars = 0
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)
                    total_chars += len(page_text)
        except Exception as e:
            print(f"❌ pdfplumber failed: {e}, trying PyPDF2...")
            return self._extract_text_pypdf2(file_path)
        
        text = "\n\n".join(text_parts)
        
        # Calculate quality score
        quality_score = self._calculate_text_quality(text, total_chars)
        
        return text, quality_score
    
    def _extract_text_pypdf2(self, file_path: str) -> tuple[str, float]:
        """
        Fallback: Extract text using PyPDF2
        """
        text_parts = []
        total_chars = 0
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)
                    total_chars += len(page_text)
        except Exception as e:
            print(f"❌ PyPDF2 also failed: {e}")
            return "", 0.0
        
        text = "\n\n".join(text_parts)
        quality_score = self._calculate_text_quality(text, total_chars)
        
        return text, quality_score
    
    def _calculate_text_quality(self, text: str, char_count: int) -> float:
        """
        Calculate text quality score (0-1)
        Low score = likely needs OCR
        """
        if not text or char_count == 0:
            return 0.0
        
        # Check for common OCR issues
        issues = 0
        
        # 1. Too many special characters (OCR artifacts)
        special_char_ratio = len(re.findall(r'[^\w\s.,!?;:\-\(\)]', text)) / max(char_count, 1)
        if special_char_ratio > 0.1:
            issues += 1
        
        # 2. Too many single characters (broken words)
        single_chars = len(re.findall(r'\b\w\b', text))
        if single_chars / max(len(text.split()), 1) > 0.3:
            issues += 1
        
        # 3. Very short text (likely failed extraction)
        if char_count < 100:
            issues += 2
        
        # 4. No spaces (corrupted extraction)
        if ' ' not in text[:1000]:
            issues += 2
        
        # Calculate score (fewer issues = higher score)
        quality_score = max(0.0, 1.0 - (issues * 0.2))
        
        return quality_score
    
    def _needs_ocr(self, text: str, quality_score: float) -> bool:
        """
        Determine if document needs OCR
        """
        # If quality is low, recommend OCR
        if quality_score < 0.5:
            return True
        
        # If text is too short for a document, likely scanned
        if len(text) < 500:
            return True
        
        return False
    
    def _needs_vision_analysis(self, file_path: str) -> bool:
        """
        Quick check if PDF has images that might need vision analysis
        This is a lightweight check - doesn't extract images yet
        """
        try:
            with pdfplumber.open(file_path) as pdf:
                # Check first few pages for images
                for page in pdf.pages[:3]:  # Only check first 3 pages
                    if len(page.images) > 0:
                        return True
            return False
        except:
            return False
    
    def _process_docx(
        self,
        file_path: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> ProcessingResult:
        """Process DOCX file"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
            text = "\n\n".join(text_parts)
            
            chunks = self._create_chunks(text, document_id, metadata)
            
            return ProcessingResult(
                success=True,
                text=text,
                chunks=chunks,
                metadata={
                    **metadata,
                    'text_length': len(text),
                    'chunk_count': len(chunks),
                    'extraction_method': 'python-docx'
                }
            )
        except ImportError:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error="python-docx not installed. Run: pip install python-docx"
            )
        except Exception as e:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error=str(e)
            )
    
    def _process_pptx(
        self,
        file_path: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> ProcessingResult:
        """Process PPTX file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
            
            text = "\n\n".join(text_parts)
            chunks = self._create_chunks(text, document_id, metadata)
            
            return ProcessingResult(
                success=True,
                text=text,
                chunks=chunks,
                metadata={
                    **metadata,
                    'text_length': len(text),
                    'chunk_count': len(chunks),
                    'slide_count': len(prs.slides),
                    'extraction_method': 'python-pptx'
                }
            )
        except ImportError:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error="python-pptx not installed. Run: pip install python-pptx"
            )
        except Exception as e:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error=str(e)
            )
    
    def _process_text(
        self,
        file_path: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> ProcessingResult:
        """Process plain text or HTML file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            chunks = self._create_chunks(text, document_id, metadata)
            
            return ProcessingResult(
                success=True,
                text=text,
                chunks=chunks,
                metadata={
                    **metadata,
                    'text_length': len(text),
                    'chunk_count': len(chunks),
                    'extraction_method': 'direct'
                }
            )
        except Exception as e:
            return ProcessingResult(
                success=False,
                text="",
                chunks=[],
                metadata={},
                error=str(e)
            )
    
    def _create_chunks(
        self,
        text: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """
        Create chunks from text with metadata
        """
        if not text or not text.strip():
            return []
        
        # Split text into chunks
        text_chunks = self.text_splitter.split_text(text)
        
        # Create chunk objects with metadata
        chunks = []
        for idx, chunk_text in enumerate(text_chunks):
            chunk = {
                'document_id': document_id,
                'chunk_index': idx,
                'content': chunk_text,
                'token_count': self.count_tokens(chunk_text),
                'metadata': {
                    **metadata,
                    'chunk_position': f"{idx + 1}/{len(text_chunks)}"
                }
            }
            chunks.append(chunk)
        
        return chunks


# Singleton instance
_processor = None

def get_processor() -> DocumentProcessor:
    """Get or create document processor instance"""
    global _processor
    if _processor is None:
        _processor = DocumentProcessor()
    return _processor
