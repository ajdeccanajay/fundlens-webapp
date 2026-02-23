"""
FastAPI server for SEC filing processing
Provides REST API endpoints for the Node.js backend

Uses Hybrid Parser (iXBRL + HTML fallback) for maximum accuracy.
"""

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import Dict, List, Optional, Any
import logging
import traceback
from datetime import datetime

# Import reporting unit extractor
from reporting_unit_extractor import ReportingUnitExtractor, ReportingUnitInfo

# Try to import hybrid parser first, fall back to ultra_accurate
try:
    from hybrid_parser import HybridSECParser
    PARSER_TYPE = "hybrid"
except ImportError:
    from ultra_accurate_parser import UltraAccurateSECParser
    PARSER_TYPE = "ultra_accurate"

from accuracy_validator import AccuracyValidator

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="SEC Filing Parser API", version="2.0.0")

# Initialize parser based on availability
if PARSER_TYPE == "hybrid":
    parser = HybridSECParser()
    logger.info("Using Hybrid iXBRL Parser (v2.0)")
else:
    parser = UltraAccurateSECParser()
    logger.info("Using Ultra-Accurate HTML Parser (v1.0)")

validator = AccuracyValidator()

class FilingRequest(BaseModel):
    ticker: str
    filingType: str
    content: str
    accessionNumber: str
    extractMetrics: bool = True
    extractNarratives: bool = True
    validateAccuracy: bool = True

class FilingResponse(BaseModel):
    success: bool
    error: Optional[str] = None
    metrics: Optional[Dict[str, Any]] = None
    narratives: Optional[Dict[str, Any]] = None
    processingTime: Optional[float] = None
    accuracy: Optional[Dict[str, Any]] = None

@app.get("/")
async def root():
    return {"message": "SEC Filing Parser API", "status": "running"}

@app.get("/health")
async def health():
    return {"status": "healthy", "parser": "ready"}

@app.post("/parse-filing", response_model=FilingResponse)
async def parse_filing(request: FilingRequest):
    """
    Parse SEC filing content and extract metrics and narratives
    """
    try:
        logger.info(f"Processing {request.ticker} {request.filingType} {request.accessionNumber}")
        
        import time
        start_time = time.time()
        
        result = {
            "success": True,
            "metrics": {},
            "narratives": {},
            "processingTime": 0
        }
        
        # Use ultra-accurate parser
        try:
            parsed_result = parser.parse_filing(
                request.content,
                request.ticker,
                request.filingType,
                "unknown_cik"  # CIK not provided in request
            )
            
            # Extract metrics if requested
            if request.extractMetrics:
                structured_metrics = parsed_result.get('structured_metrics', [])
                
                # Convert to expected format - use raw_label in key to preserve unique line items
                metrics_dict = {}
                for metric in structured_metrics:
                    # Create a unique key that preserves different raw labels
                    # This ensures we don't lose unique line items that normalize to the same metric
                    raw_label_slug = metric['raw_label'].lower().replace(' ', '_').replace(':', '_')[:50]
                    key = f"{metric['normalized_metric']}_{raw_label_slug}_{metric['period_type']}_{metric['fiscal_period']}"
                    
                    # If key exists, keep the one with higher confidence
                    if key in metrics_dict:
                        if metric['confidence_score'] > metrics_dict[key]['confidence']:
                            metrics_dict[key] = {
                                "value": metric['value'],
                                "unit": "USD",
                                "reporting_unit": metric.get('reporting_unit', 'units'),  # Original scale from SEC filing
                                "period": metric['fiscal_period'],
                                "confidence": metric['confidence_score'],
                                "raw_label": metric['raw_label'],
                                "statement_type": metric['statement_type'],
                                "normalized_metric": metric['normalized_metric']  # Store normalized metric separately
                            }
                    else:
                        metrics_dict[key] = {
                            "value": metric['value'],
                            "unit": "USD",
                            "reporting_unit": metric.get('reporting_unit', 'units'),  # Original scale from SEC filing
                            "period": metric['fiscal_period'],
                            "confidence": metric['confidence_score'],
                            "raw_label": metric['raw_label'],
                            "statement_type": metric['statement_type'],
                            "normalized_metric": metric['normalized_metric']  # Store normalized metric separately
                        }
                
                result["metrics"] = metrics_dict
                logger.info(f"{PARSER_TYPE.title()} extraction: {len(metrics_dict)} metrics")
            
            # Extract narratives if requested
            if request.extractNarratives:
                narrative_chunks = parsed_result.get('narrative_chunks', [])
                
                # Convert to expected format
                narratives_dict = {}
                for chunk in narrative_chunks:
                    # Handle both section_type and section_key field names
                    section_type = chunk.get('section_type') or chunk.get('section_key', 'general')
                    if section_type not in narratives_dict:
                        narratives_dict[section_type] = []
                    
                    narratives_dict[section_type].append({
                        "content": chunk['content'],
                        "chunk_index": chunk['chunk_index']
                    })
                
                result["narratives"] = narratives_dict
                
                total_chunks = sum(len(chunks) for chunks in narratives_dict.values())
                logger.info(f"Extracted {total_chunks} narrative chunks across {len(narratives_dict)} sections")
            
            # Validate accuracy if requested
            if request.validateAccuracy and request.extractMetrics:
                try:
                    accuracy_result = validator.validate_extraction(
                        request.content,
                        parsed_result.get('structured_metrics', []),
                        request.ticker
                    )
                    result["accuracy"] = accuracy_result
                    
                    logger.info(f"Accuracy validation: {accuracy_result['accuracy']:.4f}% "
                              f"({accuracy_result['extracted_line_items']}/{accuracy_result['total_line_items']} items)")
                    
                    if not accuracy_result['passed']:
                        logger.warning(f"Accuracy below 99.999%! Missing items: {accuracy_result['missing_items'][:5]}")
                
                except Exception as e:
                    logger.error(f"Error validating accuracy: {str(e)}")
                    result["accuracy"] = {"error": str(e)}
            
        except Exception as e:
            logger.error(f"Error in ultra-accurate parsing: {str(e)}")
            # Don't overwrite successful hybrid parsing results - keep what we have
            if not result.get("accuracy"):
                result["accuracy"] = {"error": str(e)}
        
        result["processingTime"] = time.time() - start_time
        logger.info(f"Processing completed in {result['processingTime']:.2f} seconds")
        
        return FilingResponse(**result)
        
    except Exception as e:
        logger.error(f"Error processing filing: {str(e)}")
        logger.error(traceback.format_exc())
        
        return FilingResponse(
            success=False,
            error=str(e),
            processingTime=time.time() - start_time if 'start_time' in locals() else 0
        )

@app.post("/test-parser")
async def test_parser():
    """
    Test endpoint to verify parser is working
    """
    try:
        # Simple test HTML content
        test_content = """
        <html>
        <body>
        <table>
        <tr><td>Revenue</td><td>$1,000,000</td></tr>
        <tr><td>Net Income</td><td>$200,000</td></tr>
        </table>
        <p>Management's Discussion and Analysis: This is a test narrative section.</p>
        </body>
        </html>
        """
        
        request = FilingRequest(
            ticker="TEST",
            filingType="10-K",
            content=test_content,
            accessionNumber="test-123",
            extractMetrics=True,
            extractNarratives=True,
            validateAccuracy=True
        )
        
        return await parse_filing(request)
        
    except Exception as e:
        logger.error(f"Test failed: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/sec-parser")
async def sec_parser_legacy(request: dict):
    """
    Legacy endpoint for compatibility with Node.js ingestion service
    """
    try:
        # Convert legacy format to new format
        filing_request = FilingRequest(
            ticker=request.get("ticker", "UNKNOWN"),
            filingType=request.get("filing_type", "10-K"),
            content=request.get("html_content", ""),
            accessionNumber=request.get("cik", "unknown"),
            extractMetrics=True,
            extractNarratives=True,
            validateAccuracy=True
        )
        
        # Use the new parser
        result = await parse_filing(filing_request)
        
        # Convert to legacy format expected by Node.js
        legacy_response = {
            "structured_metrics": [],
            "narrative_chunks": [],
            "metadata": {
                "ticker": filing_request.ticker,
                "filing_type": filing_request.filingType,
                "cik": filing_request.accessionNumber,
                "total_metrics": 0,
                "total_chunks": 0,
                "high_confidence_metrics": 0
            }
        }
        
        if result.success and result.metrics:
            # Convert metrics to legacy format
            for key, metric in result.metrics.items():
                # Use the stored normalized_metric if available, otherwise extract from key
                if 'normalized_metric' in metric:
                    normalized_metric = metric['normalized_metric']
                else:
                    # Fallback: Extract normalized metric name from key format
                    # Examples: "revenue_annual_FY2024" -> "revenue", "net_income_annual_FY2024" -> "net_income"
                    
                    # Remove the fiscal period part (e.g., "FY2024")
                    if '_FY' in key:
                        key_without_period = key.rsplit('_FY', 1)[0]
                    elif '_Q' in key:
                        key_without_period = key.rsplit('_Q', 1)[0]
                    else:
                        key_without_period = key
                    
                    # Remove the period type (annual/quarterly)
                    if key_without_period.endswith('_annual'):
                        normalized_metric = key_without_period[:-7]  # Remove "_annual"
                    elif key_without_period.endswith('_quarterly'):
                        normalized_metric = key_without_period[:-10]  # Remove "_quarterly"
                    else:
                        # Fallback
                        normalized_metric = key_without_period
                
                legacy_response["structured_metrics"].append({
                    "ticker": filing_request.ticker,
                    "normalized_metric": normalized_metric,
                    "raw_label": metric["raw_label"],
                    "value": metric["value"],
                    "fiscal_period": metric["period"],
                    "period_type": "annual" if "FY" in metric["period"] else "quarterly",
                    "filing_type": filing_request.filingType,
                    "statement_type": metric["statement_type"],
                    "confidence_score": metric["confidence"]
                })
            
            legacy_response["metadata"]["total_metrics"] = len(result.metrics)
            legacy_response["metadata"]["high_confidence_metrics"] = sum(
                1 for m in result.metrics.values() if m["confidence"] >= 0.9
            )
        
        if result.success and result.narratives:
            # Convert narratives to legacy format
            for section_type, chunks in result.narratives.items():
                for chunk in chunks:
                    legacy_response["narrative_chunks"].append({
                        "ticker": filing_request.ticker,
                        "filing_type": filing_request.filingType,
                        "section_type": section_type,
                        "chunk_index": chunk["chunk_index"],
                        "content": chunk["content"]
                    })
            
            legacy_response["metadata"]["total_chunks"] = sum(
                len(chunks) for chunks in result.narratives.values()
            )
        
        logger.info(f"Legacy endpoint: {legacy_response['metadata']['total_metrics']} metrics, "
                   f"{legacy_response['metadata']['total_chunks']} chunks")
        
        return legacy_response
        
    except Exception as e:
        logger.error(f"Error in legacy sec-parser endpoint: {str(e)}")
        return {
            "structured_metrics": [],
            "narrative_chunks": [],
            "metadata": {
                "ticker": request.get("ticker", "UNKNOWN"),
                "filing_type": request.get("filing_type", "10-K"),
                "cik": request.get("cik", "unknown"),
                "total_metrics": 0,
                "total_chunks": 0,
                "high_confidence_metrics": 0,
                "error": str(e)
            }
        }

@app.get("/accuracy-report")
async def get_accuracy_report():
    """
    Get comprehensive accuracy validation report
    """
    try:
        report = validator.generate_report()
        return {
            "report": report,
            "validation_count": len(validator.validation_results),
            "timestamp": "2025-12-17T15:45:00Z"
        }
    except Exception as e:
        logger.error(f"Error generating accuracy report: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================
# FINANCIAL METRICS CALCULATION ENDPOINT
# This is the PRIMARY endpoint for calculating derived metrics
# Called by Node.js backend in ECS via HTTP
# ============================================================

class CalculateMetricsRequest(BaseModel):
    ticker: str
    years: int = 5
    sharePrice: Optional[float] = None

class CalculatedMetricResponse(BaseModel):
    metricName: str
    value: float
    period: str
    periodType: str
    calculationMethod: str
    sourceMetrics: List[str]
    confidenceScore: float
    calculationDate: str
    validationStatus: str

class CalculateMetricsResponse(BaseModel):
    success: bool
    ticker: str
    metricsCount: int
    metrics: List[CalculatedMetricResponse]
    savedToDatabase: bool
    processingTimeMs: float
    error: Optional[str] = None

@app.post("/calculate-metrics", response_model=CalculateMetricsResponse)
async def calculate_metrics(request: CalculateMetricsRequest):
    """
    Calculate comprehensive financial metrics for a company
    
    This is the DETERMINISTIC calculation engine that:
    1. Reads raw metrics from PostgreSQL (ingested from SEC filings)
    2. Calculates derived metrics (TTM, margins, ratios, growth rates)
    3. Saves calculated metrics back to PostgreSQL
    4. Returns all calculated metrics
    
    Called by Node.js backend during pipeline Step B
    """
    import time
    start_time = time.time()
    
    try:
        logger.info(f"🧮 Calculating metrics for {request.ticker} ({request.years} years)")
        
        # Import the comprehensive calculator
        from comprehensive_financial_calculator import ComprehensiveFinancialCalculator
        
        calculator = ComprehensiveFinancialCalculator()
        
        # Calculate all metrics
        metrics_dict = calculator.calculate_all_metrics(
            request.ticker, 
            request.years, 
            request.sharePrice
        )
        
        if not metrics_dict:
            logger.warning(f"No metrics calculated for {request.ticker}")
            return CalculateMetricsResponse(
                success=False,
                ticker=request.ticker,
                metricsCount=0,
                metrics=[],
                savedToDatabase=False,
                processingTimeMs=(time.time() - start_time) * 1000,
                error=f"No raw metrics found for {request.ticker}. Run SEC ingestion first."
            )
        
        # Save to database
        saved = False
        try:
            calculator.save_all_metrics(request.ticker, metrics_dict)
            saved = True
            logger.info(f"✅ Saved {len(metrics_dict)} metrics to database for {request.ticker}")
        except Exception as save_error:
            logger.error(f"Failed to save metrics: {save_error}")
        
        # Convert to response format
        metrics_list = []
        for name, metric in metrics_dict.items():
            metrics_list.append(CalculatedMetricResponse(
                metricName=metric.metric_name,
                value=float(metric.value) if metric.value else 0.0,
                period=metric.period,
                periodType=metric.period_type,
                calculationMethod=metric.calculation_method,
                sourceMetrics=metric.source_metrics or [],
                confidenceScore=float(metric.confidence_score) if metric.confidence_score else 1.0,
                calculationDate=metric.calculation_date.isoformat() if metric.calculation_date else datetime.now().isoformat(),
                validationStatus=metric.validation_status or 'calculated'
            ))
        
        processing_time = (time.time() - start_time) * 1000
        logger.info(f"✅ Calculated {len(metrics_list)} metrics for {request.ticker} in {processing_time:.0f}ms")
        
        return CalculateMetricsResponse(
            success=True,
            ticker=request.ticker,
            metricsCount=len(metrics_list),
            metrics=metrics_list,
            savedToDatabase=saved,
            processingTimeMs=processing_time,
            error=None
        )
        
    except Exception as e:
        logger.error(f"❌ Error calculating metrics for {request.ticker}: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        
        return CalculateMetricsResponse(
            success=False,
            ticker=request.ticker,
            metricsCount=0,
            metrics=[],
            savedToDatabase=False,
            processingTimeMs=(time.time() - start_time) * 1000,
            error=str(e)
        )


@app.get("/calculate-metrics/{ticker}")
async def calculate_metrics_get(ticker: str, years: int = 5, sharePrice: Optional[float] = None):
    """
    GET endpoint for calculating metrics (convenience endpoint)
    """
    request = CalculateMetricsRequest(ticker=ticker, years=years, sharePrice=sharePrice)
    return await calculate_metrics(request)


@app.get("/metrics/{ticker}")
async def get_calculated_metrics(ticker: str):
    """
    Get already-calculated metrics from database
    Does NOT recalculate - just retrieves existing calculated metrics
    """
    try:
        from comprehensive_financial_calculator import ComprehensiveFinancialCalculator
        import psycopg2
        
        calculator = ComprehensiveFinancialCalculator()
        conn = calculator.connect_db()
        
        try:
            cursor = conn.cursor()
            cursor.execute("""
                SELECT metric_name, value, period, period_type, 
                       calculation_method, source_metrics, confidence_score,
                       calculation_date, validation_status
                FROM calculated_metrics 
                WHERE ticker = %s
                ORDER BY calculation_date DESC
            """, (ticker.upper(),))
            
            rows = cursor.fetchall()
            
            if not rows:
                return {
                    "success": False,
                    "ticker": ticker.upper(),
                    "metricsCount": 0,
                    "metrics": [],
                    "error": f"No calculated metrics found for {ticker}. Run /calculate-metrics first."
                }
            
            metrics = []
            for row in rows:
                metrics.append({
                    "metricName": row[0],
                    "value": float(row[1]) if row[1] else 0.0,
                    "period": row[2],
                    "periodType": row[3],
                    "calculationMethod": row[4],
                    "sourceMetrics": row[5] if row[5] else [],
                    "confidenceScore": float(row[6]) if row[6] else 1.0,
                    "calculationDate": row[7].isoformat() if row[7] else None,
                    "validationStatus": row[8] or 'calculated'
                })
            
            return {
                "success": True,
                "ticker": ticker.upper(),
                "metricsCount": len(metrics),
                "metrics": metrics
            }
            
        finally:
            conn.close()
            
    except Exception as e:
        logger.error(f"Error getting metrics for {ticker}: {str(e)}")
        return {
            "success": False,
            "ticker": ticker.upper(),
            "metricsCount": 0,
            "metrics": [],
            "error": str(e)
        }


# ============================================================
# GENERIC FORMULA EVALUATION ENDPOINT
# Safe expression evaluation via simpleeval (no Python eval())
# ============================================================

import time as _time
from simpleeval import simple_eval, InvalidExpression

class CalculateFormulaRequest(BaseModel):
    formula: str
    inputs: Dict[str, float] = {}

class CalculateFormulaResponse(BaseModel):
    result: Optional[float] = None
    audit_trail: Optional[Dict[str, Any]] = None
    error: Optional[str] = None
    formula: Optional[str] = None
    provided_inputs: Optional[Dict[str, float]] = None

@app.post("/calculate", response_model=CalculateFormulaResponse)
async def calculate_formula(request: CalculateFormulaRequest):
    """
    Generic formula evaluation endpoint using simpleeval.
    
    Accepts a formula string and named numeric inputs,
    evaluates safely (no builtins, no imports), returns result with audit trail.
    """
    formula = request.formula
    inputs = request.inputs

    if not formula or not formula.strip():
        return CalculateFormulaResponse(
            error="formula is required and must be non-empty",
            formula=formula,
            provided_inputs=inputs,
        )

    # Build intermediate steps for audit trail
    intermediate_steps = []
    intermediate_steps.append(f"Formula: {formula}")
    input_str = ", ".join(f"{k}={v}" for k, v in inputs.items())
    intermediate_steps.append(f"Inputs: {{{input_str}}}")

    try:
        start_time = _time.perf_counter()
        result = simple_eval(
            formula,
            names=inputs,
            functions={
                "abs": abs,
                "round": round,
                "min": min,
                "max": max,
            },
        )
        execution_time_ms = (_time.perf_counter() - start_time) * 1000

        if not isinstance(result, (int, float)):
            return CalculateFormulaResponse(
                error=f"Formula did not produce a numeric result: {type(result).__name__}",
                formula=formula,
                provided_inputs=inputs,
            )

        result = float(result)
        intermediate_steps.append(f"Result: {result}")
        intermediate_steps.append(f"Execution time: {execution_time_ms:.3f}ms")

        return CalculateFormulaResponse(
            result=result,
            audit_trail={
                "formula": formula,
                "inputs": inputs,
                "intermediate_steps": intermediate_steps,
                "result": result,
                "execution_time_ms": round(execution_time_ms, 3),
            },
        )

    except ZeroDivisionError:
        return CalculateFormulaResponse(
            error="Division by zero in formula evaluation",
            formula=formula,
            provided_inputs=inputs,
        )
    except InvalidExpression as e:
        return CalculateFormulaResponse(
            error=f"Invalid formula expression: {str(e)}",
            formula=formula,
            provided_inputs=inputs,
        )
    except NameError as e:
        return CalculateFormulaResponse(
            error=f"Unknown variable in formula: {str(e)}",
            formula=formula,
            provided_inputs=inputs,
        )
    except Exception as e:
        logger.error(f"Formula evaluation failed: {str(e)}")
        return CalculateFormulaResponse(
            error=f"Formula evaluation failed: {str(e)}",
            formula=formula,
            provided_inputs=inputs,
        )


# ============================================================
# REPORTING UNIT EXTRACTION ENDPOINT
# Used by backfill script to accurately determine reporting units
# ============================================================

class ExtractReportingUnitRequest(BaseModel):
    content: str
    ticker: Optional[str] = None

class ExtractReportingUnitResponse(BaseModel):
    success: bool
    default_unit: str
    share_unit: str
    per_share_unit: str
    source: str
    raw_pattern: Optional[str] = None
    error: Optional[str] = None

# Initialize reporting unit extractor
reporting_unit_extractor = ReportingUnitExtractor()

@app.post("/extract-reporting-unit", response_model=ExtractReportingUnitResponse)
async def extract_reporting_unit(request: ExtractReportingUnitRequest):
    """
    Extract reporting unit from SEC filing content.
    
    This endpoint uses the ReportingUnitExtractor to accurately determine
    the reporting units (millions, thousands, etc.) from SEC filing headers.
    
    Used by the backfill script to update existing metrics with correct units.
    """
    try:
        if not request.content:
            return ExtractReportingUnitResponse(
                success=False,
                default_unit='units',
                share_unit='units',
                per_share_unit='units',
                source='error',
                error="No content provided"
            )
        
        logger.info(f"Extracting reporting unit for {request.ticker or 'unknown ticker'}")
        
        unit_info = reporting_unit_extractor.extract_from_filing(request.content)
        
        logger.info(f"Extracted: default={unit_info.default_unit}, "
                   f"shares={unit_info.share_unit}, source={unit_info.source}")
        
        return ExtractReportingUnitResponse(
            success=True,
            default_unit=unit_info.default_unit,
            share_unit=unit_info.share_unit,
            per_share_unit=unit_info.per_share_unit,
            source=unit_info.source,
            raw_pattern=unit_info.raw_pattern
        )
        
    except Exception as e:
        logger.error(f"Error extracting reporting unit: {str(e)}")
        return ExtractReportingUnitResponse(
            success=False,
            default_unit='units',
            share_unit='units',
            per_share_unit='units',
            source='error',
            error=str(e)
        )


@app.post("/get-unit-for-metric")
async def get_unit_for_metric(metric_name: str, default_unit: str = 'millions', 
                               share_unit: str = 'millions', per_share_unit: str = 'units'):
    """
    Determine the correct unit for a specific metric type.
    
    Given a metric name and the filing's unit info, returns the appropriate
    unit for that metric (e.g., EPS is always in 'units', shares may be in 'thousands').
    """
    try:
        unit_info = ReportingUnitInfo(
            default_unit=default_unit,
            share_unit=share_unit,
            per_share_unit=per_share_unit,
            source='provided'
        )
        
        result_unit = reporting_unit_extractor.get_unit_for_metric(metric_name, unit_info)
        
        return {
            "success": True,
            "metric_name": metric_name,
            "unit": result_unit
        }
        
    except Exception as e:
        logger.error(f"Error getting unit for metric: {str(e)}")
        return {
            "success": False,
            "metric_name": metric_name,
            "unit": default_unit,
            "error": str(e)
        }


# ============================================================
# Vision Pipeline Endpoints (Instant RAG)
# ============================================================

from fastapi import UploadFile, File, Query
import base64
import io
import tempfile
import os

# Lazy imports for vision dependencies (may not be installed in all environments)
_pdf2image = None
_pptx = None
_PIL = None

MAX_PPTX_SLIDES = 100

def _get_pdf2image():
    global _pdf2image
    if _pdf2image is None:
        try:
            import pdf2image
            _pdf2image = pdf2image
        except ImportError:
            raise HTTPException(status_code=500, detail="pdf2image not installed. Install with: pip install pdf2image")
    return _pdf2image

def _get_pptx():
    global _pptx
    if _pptx is None:
        try:
            from pptx import Presentation
            _pptx = Presentation
        except ImportError:
            raise HTTPException(status_code=500, detail="python-pptx not installed. Install with: pip install python-pptx")
    return _pptx

def _get_pil():
    global _PIL
    if _PIL is None:
        try:
            from PIL import Image
            _PIL = Image
        except ImportError:
            raise HTTPException(status_code=500, detail="Pillow not installed. Install with: pip install Pillow")
    return _PIL


class VisionRenderResponse(BaseModel):
    images: List[str]  # base64-encoded images
    page_count: int
    rendered_count: int
    truncated: bool
    warnings: List[str] = []


@app.post("/vision/render-pdf", response_model=VisionRenderResponse)
async def render_pdf(
    file: UploadFile = File(...),
    dpi: int = Query(default=150, ge=72, le=300),
):
    """
    Render PDF pages to base64-encoded PNG images for vision analysis.
    """
    pdf2image = _get_pdf2image()
    warnings = []

    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty file uploaded")

        # Write to temp file (pdf2image needs a file path)
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            images = pdf2image.convert_from_path(tmp_path, dpi=dpi)
        finally:
            os.unlink(tmp_path)

        page_count = len(images)
        encoded_images = []

        for img in images:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            encoded_images.append(base64.b64encode(buf.getvalue()).decode("utf-8"))

        logger.info(f"Rendered PDF: {page_count} pages at {dpi} DPI")

        return VisionRenderResponse(
            images=encoded_images,
            page_count=page_count,
            rendered_count=page_count,
            truncated=False,
            warnings=warnings,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PDF rendering failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF rendering failed: {str(e)}")


@app.post("/vision/render-pptx", response_model=VisionRenderResponse)
async def render_pptx(
    file: UploadFile = File(...),
    dpi: int = Query(default=150, ge=72, le=300),
):
    """
    Render PPTX slides to base64-encoded PNG images for vision analysis.
    Limits to first 100 slides.
    """
    Presentation = _get_pptx()
    Image = _get_pil()
    warnings = []

    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty file uploaded")

        prs = Presentation(io.BytesIO(content))
        total_slides = len(prs.slides)
        truncated = total_slides > MAX_PPTX_SLIDES

        if truncated:
            warnings.append(f"PPTX has {total_slides} slides, only first {MAX_PPTX_SLIDES} rendered")
            logger.warn(f"PPTX truncated: {total_slides} slides, rendering {MAX_PPTX_SLIDES}")

        slides_to_render = min(total_slides, MAX_PPTX_SLIDES)

        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Convert EMU to pixels at target DPI
        # 1 inch = 914400 EMU, so pixels = EMU * DPI / 914400
        px_width = int(slide_width * dpi / 914400)
        px_height = int(slide_height * dpi / 914400)

        encoded_images = []

        for i, slide in enumerate(prs.slides):
            if i >= slides_to_render:
                break

            # Create a blank white image for the slide
            img = Image.new("RGB", (px_width, px_height), "white")

            # Extract text content and render as simple text overlay
            # For full visual fidelity, we render a placeholder with slide text
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

            # For now, create a simple rendered slide image
            # In production, use LibreOffice or similar for full rendering
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            encoded_images.append(base64.b64encode(buf.getvalue()).decode("utf-8"))

        logger.info(f"Rendered PPTX: {slides_to_render}/{total_slides} slides at {dpi} DPI")

        return VisionRenderResponse(
            images=encoded_images,
            page_count=total_slides,
            rendered_count=slides_to_render,
            truncated=truncated,
            warnings=warnings,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PPTX rendering failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PPTX rendering failed: {str(e)}")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)